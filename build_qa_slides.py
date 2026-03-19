"""
build_qa_slides.py — generates the Q&A prep deck as a .pptx.

One slide per question: hostile question as headline, sharp answer as body.
Plus a title slide and a bridging phrases slide.
Theme: corporate (dark blue — signals authority, not defensiveness).
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colors — corporate theme
# ---------------------------------------------------------------------------
BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
SLIDE   = RGBColor(0x1E, 0x29, 0x3B)   # slate
ACCENT  = RGBColor(0x0E, 0xA5, 0xE9)   # sky blue
TEXT    = RGBColor(0xF1, 0xF5, 0xF9)   # near white
MUTED   = RGBColor(0x94, 0xA3, 0xB8)   # slate gray
DIVIDER = RGBColor(0x33, 0x41, 0x55)   # dark border
RED     = RGBColor(0xEF, 0x44, 0x44)   # question highlight

# ---------------------------------------------------------------------------
# Q&A content
# ---------------------------------------------------------------------------

QA = [
    (
        "\"How is this different from just asking ChatGPT to write my slides?\"",
        "ChatGPT gives you text. You still paste it into PowerPoint, format it, and structure it.\n\n"
        "This outputs a native .pptx — layout, speaker notes, theming — from one command.\n\n"
        "The difference: a sous chef hands you instructions. This hands you dinner."
    ),
    (
        "\"What happens when the AI gets the content wrong?\"",
        "The same thing that happens when a junior analyst gets it wrong — you catch it in review.\n\n"
        "Except you're reviewing a 60-second first draft, not a four-hour one.\n\n"
        "Your attention shifts from building slides to catching errors. That's the right job."
    ),
    (
        "\"We're regulated. Can we send our topics to an external API?\"",
        "What goes to the API is one sentence: your topic. Not your data. Not your documents.\n\n"
        "'Q3 board update' is the level of sensitivity we're talking about.\n\n"
        "For air-gapped environments: swap the endpoint to a self-hosted model. One-line change."
    ),
    (
        "\"What's the quality ceiling? Can I actually send this to a client?\"",
        "I'll answer that with a question: what's the quality ceiling of a first draft from your best analyst?\n\n"
        "That analyst still needs your feedback, your edits, your judgment. So does this.\n\n"
        "What the tool produces is a structured argument — the 80% that doesn't require you. "
        "You spend your time on the 20% that does. You're still the author."
    ),
    (
        "\"What about complex, multi-stakeholder presentations?\"",
        "Two options.\n\n"
        "One: use the tool for structure-heavy, low-judgment sections — agenda, context, appendix.\n\n"
        "Two: use --slides to feed your own JSON. Your data, your sections. It handles layout and notes.\n\n"
        "You're the architect. It's the builder."
    ),
    (
        "\"What does this cost at scale for a team of 50?\"",
        "Approximately $0.03–$0.05 per presentation at current API pricing.\n\n"
        "For 50 people, one deck per day: under $1,000 a year.\n\n"
        "Compare that to what you're currently spending on that output in senior labor.\n\n"
        "The math takes thirty seconds."
    ),
    (
        "\"Could a competitor clone this in a weekend?\"",
        "Let me separate two questions asked as one.\n\n"
        "Can someone clone the code? Yes, in a weekend. That's happened. We know.\n\n"
        "Can someone clone the outcome? No. The slide schema, prompt architecture, speaker note logic — those took iteration. You can copy the output. You can't copy the iteration.\n\n"
        "The real moat is workflow. The moment a team's first draft is always a generated draft, switching cost compounds. Google Slides is 15 years old. PowerPoint is 40. Neither generates the first draft. That's the gap."
    ),
    (
        "\"Why would I trust AI to represent my thinking in a boardroom?\"",
        "You wouldn't — and you shouldn't. The tool doesn't ask you to.\n\n"
        "It generates a frame. You fill it with your thinking.\n\n"
        "The question is whether you'd rather start from a blank page or a structured draft.\n\n"
        "Nobody trusts a blank page. You trust your ability to improve on a draft."
    ),
    (
        "\"What if Anthropic changes pricing or discontinues the API?\"",
        "The tool makes one API call, abstracted behind a single function: generate_content().\n\n"
        "Swapping the model is a one-line change. OpenAI, Gemini, self-hosted Llama — all compatible.\n\n"
        "You're not locked into Anthropic. You're locked into a function signature.\n\n"
        "That's a meaningful distinction."
    ),
    (
        "\"Why is this open source? What's the business model?\"",
        "Three honest answers, depending on who's asking.\n\n"
        "Developer evaluating whether to build on this: free and open source. Cost of a presentation is three cents in API calls.\n\n"
        "Company evaluating adoption at scale: roughly $500–$1,000/year for a team of 50. No license fee. No per-seat pricing.\n\n"
        "If you're asking whether this becomes a business: yes. Open-source core stays free. Revenue comes from the layer above — hosted version, team management, brand theme marketplace, enterprise SSO. The pattern is Hashicorp, Grafana, GitLab. Give the engine away. Sell the dashboard."
    ),
]

BRIDGES = [
    (
        "Bridge 1 — The Reframe",
        "Use when: a question attacks a specific limitation or edge case.",
        "\"That's exactly the right constraint to put on it —\nand here's how I'd think about it in that context...\""
    ),
    (
        "Bridge 2 — The Acknowledge-and-Advance",
        "Use when: the concern is legitimate and you don't have a perfect answer.",
        "\"Fair point — and I won't pretend it's fully solved.\nWhat I will tell you is that the trade-off looks like this...\""
    ),
    (
        "Bridge 3 — The Return",
        "Use when: the question is pulling the room into a rabbit hole.",
        "\"I want to come back to that — but let me anchor it\nto the core question first, because I think it changes the answer...\""
    ),
]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text, left, top, width, height,
            size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return tb


def divider(slide, top, color=DIVIDER):
    d = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(9), Inches(0.03))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()


def accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide)

    stripe = slide.shapes.add_shape(1, 0, Inches(2.4), Inches(10), Inches(0.8))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    textbox(slide, "Q&A PREPARATION", 0.6, 1.2, 9, 0.6,
            size=14, bold=True, color=MUTED, align=PP_ALIGN.LEFT)
    textbox(slide, "10 Hardest Questions.\nSharp Answers.", 0.6, 1.7, 9, 1.4,
            size=36, bold=True, color=TEXT, align=PP_ALIGN.LEFT)
    textbox(slide, "AI Presentation Generator  ·  For: Technical decision-makers, founders, CTOs",
            0.6, 3.5, 9, 0.5, size=14, color=MUTED, align=PP_ALIGN.LEFT)
    textbox(slide, "Never defend. Reframe, redirect, and close.",
            0.6, 4.1, 9, 0.5, size=13, color=MUTED, align=PP_ALIGN.LEFT)


def qa_slide(prs, number, question, answer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    # Question number badge
    badge = slide.shapes.add_shape(1, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    textbox(slide, str(number), 0.52, 0.2, 0.46, 0.45,
            size=16, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    # Question
    textbox(slide, question, 1.15, 0.18, 8.3, 0.9,
            size=17, bold=True, color=RED, align=PP_ALIGN.LEFT)

    divider(slide, 1.22)

    # Answer — handle multi-paragraph
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(9), Inches(3.9))
    tf = tb.text_frame
    tf.word_wrap = True
    paragraphs = answer.strip().split("\n\n")
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4) if i > 0 else Pt(0)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = para.replace("\n", " ")
        run.font.size = Pt(16)
        run.font.name = "Calibri"
        run.font.color.rgb = TEXT

    # Slide number
    textbox(slide, f"{number} / 10", 8.8, 5.2, 1, 0.35,
            size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def bridge_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide)

    textbox(slide, "3 BRIDGING PHRASES", 0.5, 0.15, 9, 0.5,
            size=13, bold=True, color=MUTED)
    textbox(slide, "When the room tries to derail the narrative.", 0.5, 0.6, 9, 0.45,
            size=18, color=TEXT)
    divider(slide, 1.15)

    tops = [1.3, 2.65, 4.0]
    for (title, usage, phrase), top in zip(BRIDGES, tops):
        # Title
        textbox(slide, title, 0.5, top, 9, 0.35, size=13, bold=True, color=ACCENT)
        # Usage
        textbox(slide, usage, 0.5, top + 0.32, 9, 0.3, size=11, color=MUTED)
        # Phrase box background
        box = slide.shapes.add_shape(1, Inches(0.5), Inches(top + 0.62), Inches(9), Inches(0.52))
        box.fill.solid()
        box.fill.fore_color.rgb = SLIDE
        box.line.color.rgb = DIVIDER
        textbox(slide, phrase, 0.65, top + 0.65, 8.7, 0.46,
                size=12, bold=True, color=TEXT)


def golden_rule_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide, RED)

    textbox(slide, "THE GOLDEN RULE OF Q&A", 0.6, 0.25, 9, 0.5,
            size=13, bold=True, color=MUTED)
    textbox(slide,
            "Never answer the question they asked.\nAnswer the question they meant.",
            0.6, 0.9, 9, 1.2, size=32, bold=True, color=TEXT)
    divider(slide, 2.35, ACCENT)
    textbox(slide,
            "Every hostile question has a fear underneath it.\n"
            "Find the fear. Address it directly. Close toward your narrative.\n\n"
            "The presenter who wins Q&A isn't the one with the best answers.\n"
            "It's the one who stays on offense while the room thinks they're playing defense.",
            0.6, 2.55, 9, 2.5, size=16, color=TEXT)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    title_slide(prs)

    for i, (question, answer) in enumerate(QA, start=1):
        qa_slide(prs, i, question, answer)

    bridge_slide(prs)
    golden_rule_slide(prs)

    out = "QA_Prep_Slides.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({1 + len(QA) + 2} slides)")


if __name__ == "__main__":
    build()
