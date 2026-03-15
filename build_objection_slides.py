"""
build_objection_slides.py — generates Objection_Handling_Slides.pptx

Six objection slides + master close + delivery principles.
Content mirrors OBJECTION_SLIDES.md exactly.
Theme: corporate dark (matches QA deck).
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
SLIDE   = RGBColor(0x1E, 0x29, 0x3B)   # slate
ACCENT  = RGBColor(0x0E, 0xA5, 0xE9)   # sky blue
TEXT    = RGBColor(0xF1, 0xF5, 0xF9)   # near white
MUTED   = RGBColor(0x94, 0xA3, 0xB8)   # slate gray
DIVIDER = RGBColor(0x33, 0x41, 0x55)   # dark border
GREEN   = RGBColor(0x22, 0xC5, 0x5E)   # concede signal
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # reframe signal
RED     = RGBColor(0xEF, 0x44, 0x44)   # close signal

# ---------------------------------------------------------------------------
# Objection content
# ---------------------------------------------------------------------------

OBJECTIONS = [
    {
        "number": 1,
        "headline": "\"The output won't be good enough to use directly.\"",
        "reframe_label": "CONCEDE FIRST",
        "reframe": "\"You're right. It's not finished. It's 80% done in 60 seconds.\"",
        "bullets": [
            "The tool generates structure, argument, and speaker notes — not pixel-perfect design",
            "Your job shrinks from \"build a deck\" to \"review and refine\" — 4 hours becomes 20 minutes",
            "Every deck you've ever approved went through revisions — this just moves you to revision 1 instantly",
        ],
        "quote": (
            "\"You don't judge a sous chef by the fact that they didn't plate it.\n"
            "You judge them by the fact that dinner is ready and you didn't have to cook.\""
        ),
        "note": (
            "Show the generated deck. Point to one slide that needed a tweak.\n"
            "Say: \"This took me 45 seconds to fix. The alternative was starting from scratch.\""
        ),
    },
    {
        "number": 2,
        "headline": "\"Our data is confidential. What goes to the API?\"",
        "reframe_label": "THE FACT",
        "reframe": "\"Only the topic goes to Claude. Nothing else.\"",
        "bullets": [
            "The tool sends one thing to the API: your topic sentence",
            "No documents. No internal data. No files. One sentence in, structured slides out.",
            "For air-gapped environments: swap the API call for a self-hosted model — the architecture is identical",
        ],
        "quote": (
            "\"You've typed topics into Google. This is the same surface area of exposure.\n"
            "The difference is Google has your entire search history.\""
        ),
        "note": (
            "If the room is enterprise-security conscious, point to the JSON schema output.\n"
            "The API call is auditable and transparent — it contains exactly what you type as the topic, nothing more."
        ),
    },
    {
        "number": 3,
        "headline": "\"We already have Canva, Google Slides, and templates.\"",
        "reframe_label": "THE DISTINCTION",
        "reframe": "\"Those tools solve the wrong problem.\"",
        "bullets": [
            "Templates give you empty boxes — you still write every word",
            "Canva makes formatting faster — the content problem remains untouched",
            "Google Slides is a canvas — this is a co-writer",
        ],
        "quote": (
            "\"Canva addresses the tool. This addresses the task.\n"
            "The task is: what do I say, in what order, and why does it matter?\n"
            "Canva never touches that. Neither does Slides. Neither does PowerPoint.\n"
            "This does.\""
        ),
        "note": (
            "Don't cite a formatting percentage stat — it's unverifiable and a skeptic will challenge it.\n"
            "Instead: ask them. \"How long did your last deck take, start to finish?\" Wait for the answer.\n"
            "Then say: \"This does that in 60 seconds. Want to watch again?\" Their own number is more powerful."
        ),
    },
    {
        "number": 4,
        "headline": "\"We'd be too dependent on Anthropic. What if costs spike?\"",
        "reframe_label": "THE REFRAME",
        "reframe": "\"You're already dependent on infrastructure you don't own.\"",
        "bullets": [
            "Your slides are in Google Drive — dependent on Google",
            "Your video calls run on AWS — dependent on Amazon",
            "Your email runs on Microsoft — dependent on Microsoft",
            "Claude API uptime: 99.9%. Cost per presentation: under $0.05 at current pricing.",
        ],
        "quote": (
            "\"The question isn't whether to depend on infrastructure.\n"
            "The question is whether the dependency is worth the return.\n"
            "At $0.05 a deck, the math is not complicated.\""
        ),
        "note": (
            "If they push on vendor lock-in: the JSON schema output is portable.\n"
            "Point --slides at any JSON file from any source. Swap the model in one line of code."
        ),
    },
    {
        "number": 5,
        "headline": "\"Our presentations require specialized knowledge AI can't replicate.\"",
        "reframe_label": "THE FLIP",
        "reframe": "\"Then you give it the knowledge. It handles everything else.\"",
        "bullets": [
            "Use --slides to feed your own JSON — your data, your structure, your domain expertise",
            "The tool generates the frame; your team fills the substance",
            "Claude was trained on a broad corpus of business writing, frameworks, and case studies",
        ],
        "quote": (
            "\"You don't ask a consultant to know your business on day one.\n"
            "You brief them. They build the deck.\n"
            "This is the same relationship — without the invoice.\""
        ),
        "note": (
            "Demo the --slides flag if time allows. Show a JSON file with pre-populated data going in,\n"
            "a fully formatted deck coming out. The tool becomes infrastructure, not a black box."
        ),
    },
    {
        "number": 6,
        "headline": "\"It won't match our brand guidelines.\"",
        "reframe_label": "THE CLOSE",
        "reframe": "\"10 lines of Python. Your brand is in.\"",
        "bullets": [
            "Themes are a dictionary in the source code — colors, fonts, layout in one place",
            "Adding a custom theme takes under 10 minutes and is immediately available via --theme",
            "The HTML output uses CSS custom properties — any designer can extend it without touching Python",
        ],
        "quote": (
            "\"The question isn't whether this can match your brand.\n"
            "The question is whether you want it to.\n"
            "If yes, it's a Friday afternoon project.\""
        ),
        "note": (
            "Show the THEMES dict. Count the lines out loud if the room is technical.\n"
            "If non-technical: \"Your designer tells us the hex codes. We're done before lunch.\""
        ),
    },
]

MASTER_CLOSE = (
    "\"That's a real concern. Here's how I'd think about it —\"\n\n"
    "[Answer in one sentence.]\n\n"
    "\"But here's what I'd ask you to hold onto: every concern you've raised is a "
    "configuration problem, not a fundamental problem. The fundamental question is whether "
    "you want to spend four hours building slides or twenty minutes reviewing them. "
    "Everything else is details.\"\n\n"
    "[Stop. Do not add anything.]"
)

DELIVERY = [
    ("CONCEDE", "Say \"You're right\" or \"That's a fair concern\" before anything else. It disarms."),
    ("FLIP", "Reframe the objection as a feature, a choice, or a non-issue with evidence."),
    ("CLOSE", "End every objection slide with a forcing function. Not a summary. A question or a consequence."),
]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, text, left, top, width, height,
            size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return tb


def multiline_textbox(slide, paragraphs, left, top, width, height, size=13, color=TEXT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3) if i > 0 else Pt(0)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = para
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = color
    return tb


def accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def divider(slide, top, color=DIVIDER):
    d = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(9), Inches(0.03))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()


def badge(slide, text, left, top, width, height, bg_color=ACCENT):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()
    textbox(slide, text, left + 0.03, top + 0.03, width - 0.06, height - 0.06,
            size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)


def quote_box(slide, text, top, height=0.9):
    box = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(9), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = SLIDE
    box.line.color.rgb = ACCENT
    textbox(slide, text, 0.65, top + 0.08, 8.7, height - 0.16,
            size=12, bold=True, color=TEXT)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide)

    stripe = slide.shapes.add_shape(1, 0, Inches(2.4), Inches(10), Inches(0.7))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    textbox(slide, "OBJECTION HANDLING", 0.6, 1.2, 9, 0.5,
            size=13, bold=True, color=MUTED)
    textbox(slide, "6 Objections.\nSharp Answers.", 0.6, 1.65, 9, 1.3,
            size=36, bold=True, color=TEXT)
    textbox(slide, "AI Presentation Generator  ·  For: Technical decision-makers, founders, CTOs",
            0.6, 3.35, 9, 0.4, size=13, color=MUTED)
    textbox(slide, "Every concern the room will have. Answered before they ask.",
            0.6, 3.85, 9, 0.4, size=12, color=MUTED)
    textbox(slide, "Principle: concede first, flip second, close with a forcing function.",
            0.6, 4.25, 9, 0.4, size=11, color=MUTED)


def objection_slide(prs, obj):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    # Number badge
    badge(slide, str(obj["number"]), 0.5, 0.15, 0.45, 0.42)

    # Headline
    textbox(slide, obj["headline"], 1.08, 0.13, 8.4, 0.52,
            size=16, bold=True, color=RED)

    divider(slide, 0.72)

    # Reframe label + text
    badge(slide, obj["reframe_label"], 0.5, 0.82, 1.6, 0.28, bg_color=AMBER)
    textbox(slide, obj["reframe"], 2.22, 0.82, 7.2, 0.32,
            size=13, bold=True, color=TEXT)

    divider(slide, 1.18)

    # Bullets
    bullet_top = 1.28
    for i, bullet in enumerate(obj["bullets"]):
        dot = slide.shapes.add_shape(1, Inches(0.5), Inches(bullet_top + i * 0.52 + 0.08),
                                     Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        textbox(slide, bullet, 0.72, bullet_top + i * 0.52, 8.7, 0.48,
                size=12, color=TEXT)

    bullet_bottom = bullet_top + len(obj["bullets"]) * 0.52

    # Quote / reframe box
    lines = obj["quote"].split("\n")
    q_height = max(0.75, 0.28 + len(lines) * 0.22)
    quote_box(slide, obj["quote"], bullet_bottom + 0.05, q_height)

    # Slide counter
    total = len(OBJECTIONS)
    textbox(slide, f"{obj['number']} / {total}", 8.8, 5.2, 1, 0.3,
            size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def master_close_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide, GREEN)

    textbox(slide, "THE OBJECTION MASTER CLOSE", 0.5, 0.15, 9, 0.4,
            size=12, bold=True, color=MUTED)
    textbox(slide, "Use this if someone raises a concern not covered above.", 0.5, 0.55, 9, 0.35,
            size=15, color=TEXT)
    divider(slide, 0.97)

    box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(3.6))
    box.fill.solid()
    box.fill.fore_color.rgb = SLIDE
    box.line.color.rgb = DIVIDER

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(1.2), Inches(8.7), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    paras = MASTER_CLOSE.split("\n\n")
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6) if i > 0 else Pt(0)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = para.replace("\n", " ")
        run.font.size = Pt(13)
        run.font.name = "Calibri"
        run.font.color.rgb = TEXT
        if para.startswith('"') or para.startswith('['):
            run.font.bold = True


def delivery_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide)

    textbox(slide, "DELIVERY PRINCIPLES", 0.5, 0.15, 9, 0.4,
            size=12, bold=True, color=MUTED)
    textbox(slide, "The concede — flip — close structure.", 0.5, 0.55, 9, 0.35,
            size=16, color=TEXT)
    divider(slide, 0.97)

    tops = [1.1, 2.4, 3.7]
    colors = [GREEN, AMBER, RED]
    for (label, desc), top, col in zip(DELIVERY, tops, colors):
        badge(slide, label, 0.5, top, 1.4, 0.38, bg_color=col)
        textbox(slide, desc, 2.05, top + 0.04, 7.4, 0.55, size=13, color=TEXT)
        if top < 3.7:
            divider(slide, top + 1.15)

    divider(slide, 4.82)
    textbox(slide,
            "The one rule: Never argue. Concede, reframe, close. "
            "The audience came to be convinced, not defeated.",
            0.5, 4.95, 9, 0.55, size=12, bold=True, color=MUTED)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    title_slide(prs)
    for obj in OBJECTIONS:
        objection_slide(prs, obj)
    master_close_slide(prs)
    delivery_slide(prs)

    out = "Objection_Handling_Slides.pptx"
    prs.save(out)
    total = 1 + len(OBJECTIONS) + 2
    print(f"Saved: {out}  ({total} slides)")


if __name__ == "__main__":
    build()
