"""
build_closing_slide.py — generates the single closing slide as a .pptx.

Design spec from CLOSING_SLIDE.md:
  - White background
  - Headline: "Every day you wait, you pay for it."
  - Three metric lines: cost today / cost with tool / payback period
  - One command at the bottom
  - No logo. No decoration. The numbers are the design.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x11, 0x18, 0x27)
RED    = RGBColor(0xDC, 0x26, 0x26)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
CODE   = RGBColor(0x1E, 0x29, 0x3B)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def textbox(slide, text, left, top, width, height,
            size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return tb


def metric_row(slide, label, value, top):
    """One metric line: large value on left, label on right."""
    textbox(slide, value, 0.6, top, 3.5, 0.65, size=36, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    textbox(slide, label, 4.0, top + 0.1, 5.8, 0.55, size=18, color=GRAY, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # White background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Thin red accent bar at top
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    # Headline
    textbox(slide, "Every day you wait, you pay for it.",
            0.6, 0.25, 9, 0.8, size=28, bold=True, color=BLACK)

    # Divider
    div = slide.shapes.add_shape(1, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    div.line.fill.background()

    # Three metric rows
    metric_row(slide, "what one presentation costs your organization today", "$[X]",      1.4)
    metric_row(slide, "what it costs with this tool",                        "60 sec",    2.2)
    metric_row(slide, "until it pays for itself completely",                 "[X] days",  3.0)

    # Divider before command
    div2 = slide.shapes.add_shape(1, Inches(0.6), Inches(3.85), Inches(8.8), Inches(0.02))
    div2.fill.solid()
    div2.fill.fore_color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    div2.line.fill.background()

    # Command block background
    cmd_bg = slide.shapes.add_shape(1, Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.9))
    cmd_bg.fill.solid()
    cmd_bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    cmd_bg.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)

    # Command text
    textbox(slide, 'python generate.py "Your next meeting"',
            0.85, 4.1, 8.3, 0.65, size=20, bold=True, color=CODE,
            align=PP_ALIGN.LEFT, italic=False)

    # Speaker note
    slide.notes_slide.notes_text_frame.text = (
        "Point at the command. Do not speak.\n"
        "Hold for five seconds.\n"
        "Then: 'Type a topic. Hit enter. Read what comes back.'\n"
        "Sit down. Say nothing else. The silence is the close."
    )

    out = "Closing_Slide.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
