"""
test_generate.py - Test suite for generate.py

Run with:
    pytest test_generate.py -v

No ANTHROPIC_API_KEY required — all Claude API calls are mocked.
"""

import json
import os
import tempfile
import time
import unittest
from unittest.mock import MagicMock, patch

from pptx.dml.color import RGBColor

import generate
from generate import (
    THEMES,
    DEFAULT_THEME,
    TOPIC_MAX_LENGTH,
    SLIDE_SCHEMA,
    _css,
    _slide_html,
    _pptx_textbox,
    _pptx_set_bg,
    build_pptx,
    build_html,
    generate_content,
    validate_topic,
    validate_output_path,
    _check_rate_limit,
)


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

SAMPLE_DATA = {
    "title": "Test Presentation",
    "subtitle": "A test subtitle",
    "slides": [
        {"type": "content", "title": "Agenda", "bullets": ["Point A", "Point B", "Point C"], "notes": "Intro notes"},
        {"type": "section", "title": "Part One", "bullets": [], "notes": ""},
        {"type": "stat", "title": "Key Metric", "bullets": [], "stat": "42%", "stat_label": "improvement", "notes": ""},
        {"type": "quote", "title": "Insight", "bullets": [], "quote": "Great quote here", "attribution": "Someone", "notes": ""},
        {"type": "content", "title": "Conclusion", "bullets": ["Wrap up", "Next steps"], "notes": "Final notes"},
    ],
}


# ---------------------------------------------------------------------------
# Theme tests
# ---------------------------------------------------------------------------

class TestThemes(unittest.TestCase):

    def test_all_themes_present(self):
        for name in ("dark", "light", "corporate", "executive"):
            self.assertIn(name, THEMES)

    def test_default_theme_exists(self):
        self.assertIn(DEFAULT_THEME, THEMES)

    def test_each_theme_has_pptx_keys(self):
        pptx_keys = {"bg", "accent", "text", "subtext", "divider"}
        for name, theme in THEMES.items():
            with self.subTest(theme=name):
                self.assertTrue(pptx_keys.issubset(theme.keys()))

    def test_each_theme_has_html_keys(self):
        html_keys = {"background", "slide_bg", "primary", "secondary",
                     "html_text", "muted", "border", "code_bg", "font_family"}
        for name, theme in THEMES.items():
            with self.subTest(theme=name):
                self.assertTrue(html_keys.issubset(theme.keys()))

    def test_pptx_colors_are_rgbcolor(self):
        for name, theme in THEMES.items():
            for key in ("bg", "accent", "text", "subtext", "divider"):
                with self.subTest(theme=name, key=key):
                    self.assertIsInstance(theme[key], RGBColor)

    def test_html_colors_are_hex_strings(self):
        for name, theme in THEMES.items():
            for key in ("background", "slide_bg", "primary", "secondary",
                        "html_text", "muted", "border", "code_bg"):
                with self.subTest(theme=name, key=key):
                    val = theme[key]
                    self.assertIsInstance(val, str)
                    self.assertTrue(val.startswith("#"), f"{key}={val!r} should start with #")


# ---------------------------------------------------------------------------
# Slide schema tests
# ---------------------------------------------------------------------------

class TestSlideSchema(unittest.TestCase):

    def test_schema_requires_title_and_slides(self):
        for field in ("title", "subtitle", "slides"):
            self.assertIn(field, SLIDE_SCHEMA["required"])

    def test_slide_item_required_fields(self):
        item = SLIDE_SCHEMA["properties"]["slides"]["items"]
        for field in ("type", "title", "bullets", "notes"):
            self.assertIn(field, item["required"])

    def test_slide_type_enum(self):
        item = SLIDE_SCHEMA["properties"]["slides"]["items"]
        enum = item["properties"]["type"]["enum"]
        for t in ("content", "section", "quote", "stat"):
            self.assertIn(t, enum)

    def test_no_additional_properties_on_root(self):
        self.assertFalse(SLIDE_SCHEMA.get("additionalProperties", True))

    def test_no_additional_properties_on_slide_item(self):
        item = SLIDE_SCHEMA["properties"]["slides"]["items"]
        self.assertFalse(item.get("additionalProperties", True))


# ---------------------------------------------------------------------------
# HTML builder tests
# ---------------------------------------------------------------------------

class TestCss(unittest.TestCase):

    def test_css_contains_theme_colors(self):
        theme = THEMES["dark"]
        css = _css(theme)
        self.assertIn(theme["background"], css)
        self.assertIn(theme["primary"], css)
        self.assertIn(theme["font_family"], css)

    def test_css_defines_css_variables(self):
        css = _css(THEMES["light"])
        for var in ("--bg", "--primary", "--text", "--muted", "--border"):
            self.assertIn(var, css)


class TestSlideHtml(unittest.TestCase):

    def test_first_slide_gets_title_class(self):
        slide = {"title": "My Title", "bullets": []}
        html = _slide_html(slide, 1, 5)
        self.assertIn("slide-type-title", html)
        self.assertIn("My Title", html)

    def test_last_slide_gets_closing_class(self):
        slide = {"title": "The End", "bullets": []}
        html = _slide_html(slide, 5, 5)
        self.assertIn("slide-type-closing", html)

    def test_middle_slide_renders_bullets(self):
        slide = {"type": "content", "title": "Section", "bullets": ["Point 1", "Point 2"]}
        html = _slide_html(slide, 2, 5)
        self.assertIn("Point 1", html)
        self.assertIn("Point 2", html)
        self.assertNotIn("slide-type-title", html)
        self.assertNotIn("slide-type-closing", html)

    def test_slide_number_in_output(self):
        slide = {"type": "content", "title": "X", "bullets": []}
        html = _slide_html(slide, 3, 7)
        self.assertIn("3", html)
        self.assertIn("7", html)

    def test_empty_bullets_renders_empty_list(self):
        slide = {"type": "content", "title": "No bullets", "bullets": []}
        html = _slide_html(slide, 2, 4)
        self.assertIn("<ul></ul>", html)

    def test_section_slide_type(self):
        slide = {"type": "section", "title": "Part Two", "bullets": []}
        html = _slide_html(slide, 2, 5)
        self.assertIn("slide-type-section", html)
        self.assertIn("Part Two", html)

    def test_quote_slide_type(self):
        slide = {"type": "quote", "title": "Q", "bullets": [],
                 "quote": "Be the change", "attribution": "Gandhi"}
        html = _slide_html(slide, 2, 5)
        self.assertIn("slide-type-quote", html)
        self.assertIn("Be the change", html)
        self.assertIn("Gandhi", html)

    def test_stat_slide_type(self):
        slide = {"type": "stat", "title": "Growth", "bullets": [],
                 "stat": "73%", "stat_label": "year-over-year"}
        html = _slide_html(slide, 2, 5)
        self.assertIn("slide-type-stat", html)
        self.assertIn("73%", html)
        self.assertIn("year-over-year", html)


class TestBuildHtml(unittest.TestCase):

    def _build(self, theme_name="dark"):
        theme = THEMES[theme_name]
        with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w") as f:
            path = f.name
        try:
            build_html(SAMPLE_DATA, theme, path)
            with open(path, encoding="utf-8") as f:
                return f.read()
        finally:
            os.unlink(path)

    def test_output_is_valid_html_skeleton(self):
        html = self._build()
        self.assertIn("<!DOCTYPE html>", html)
        self.assertIn("<html", html)
        self.assertIn("</html>", html)

    def test_title_appears_in_output(self):
        html = self._build()
        self.assertIn("Test Presentation", html)

    def test_all_slide_titles_present(self):
        html = self._build()
        # quote slides render the quote text, not the title
        for slide in SAMPLE_DATA["slides"]:
            if slide["type"] == "quote":
                self.assertIn(slide.get("quote", ""), html)
            else:
                self.assertIn(slide["title"], html)

    def test_bullets_present(self):
        html = self._build()
        self.assertIn("Point A", html)
        self.assertIn("Point B", html)

    def test_slide_count_matches(self):
        # title slide prepended + 5 content slides = 6 total
        html = self._build()
        self.assertEqual(html.count('<section class="slide'), 6)

    def test_all_themes_produce_output(self):
        for theme_name in THEMES:
            with self.subTest(theme=theme_name):
                html = self._build(theme_name)
                self.assertGreater(len(html), 500)

    def test_theme_name_badge_in_output(self):
        html = self._build("corporate")
        self.assertIn("Corporate", html)


# ---------------------------------------------------------------------------
# PPTX builder tests
# ---------------------------------------------------------------------------

class TestBuildPptx(unittest.TestCase):

    def _build(self, theme_name="dark"):
        from pptx import Presentation
        theme = THEMES[theme_name]
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            build_pptx(SAMPLE_DATA, theme, path)
            prs = Presentation(path)
            return prs
        finally:
            os.unlink(path)

    def test_slide_count(self):
        # 1 title slide + 5 content/section/stat/quote slides
        prs = self._build()
        self.assertEqual(len(prs.slides), 6)

    def test_all_themes_produce_valid_pptx(self):
        from pptx import Presentation
        for theme_name in THEMES:
            with self.subTest(theme=theme_name):
                prs = self._build(theme_name)
                self.assertGreater(len(prs.slides), 0)

    def test_slide_dimensions(self):
        from pptx.util import Inches
        prs = self._build()
        self.assertAlmostEqual(prs.slide_width.inches, 10.0, places=1)
        self.assertAlmostEqual(prs.slide_height.inches, 5.625, places=2)

    def test_output_file_is_written(self):
        theme = THEMES["dark"]
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            build_pptx(SAMPLE_DATA, theme, path)
            self.assertGreater(os.path.getsize(path), 0)
        finally:
            os.unlink(path)

    def test_notes_added_to_slide(self):
        prs = self._build()
        # slide index 1 = first content slide (Agenda, type=content), notes="Intro notes"
        notes = prs.slides[1].notes_slide.notes_text_frame.text
        self.assertEqual(notes, "Intro notes")


# ---------------------------------------------------------------------------
# generate_content tests (mocked API)
# ---------------------------------------------------------------------------

class TestGenerateContent(unittest.TestCase):

    def _make_mock_response(self, data: dict):
        text_block = MagicMock()
        text_block.type = "text"
        text_block.text = json.dumps(data)
        response = MagicMock()
        response.content = [text_block]
        return response

    @patch("generate.anthropic.Anthropic")
    def test_returns_parsed_dict(self, mock_anthropic_cls):
        mock_client = MagicMock()
        mock_anthropic_cls.return_value = mock_client
        mock_client.messages.create.return_value = self._make_mock_response(SAMPLE_DATA)

        result = generate_content("Test Topic")
        self.assertEqual(result["title"], "Test Presentation")
        self.assertEqual(len(result["slides"]), 5)

    @patch("generate.anthropic.Anthropic")
    def test_calls_correct_model(self, mock_anthropic_cls):
        mock_client = MagicMock()
        mock_anthropic_cls.return_value = mock_client
        mock_client.messages.create.return_value = self._make_mock_response(SAMPLE_DATA)

        generate_content("Some Topic")
        call_kwargs = mock_client.messages.create.call_args[1]
        self.assertEqual(call_kwargs["model"], "claude-opus-4-6")

    @patch("generate.anthropic.Anthropic")
    def test_thinking_enabled(self, mock_anthropic_cls):
        mock_client = MagicMock()
        mock_anthropic_cls.return_value = mock_client
        mock_client.messages.create.return_value = self._make_mock_response(SAMPLE_DATA)

        generate_content("Some Topic")
        call_kwargs = mock_client.messages.create.call_args[1]
        self.assertEqual(call_kwargs["thinking"]["type"], "adaptive")

    @patch("generate.anthropic.Anthropic")
    def test_topic_in_prompt(self, mock_anthropic_cls):
        mock_client = MagicMock()
        mock_anthropic_cls.return_value = mock_client
        mock_client.messages.create.return_value = self._make_mock_response(SAMPLE_DATA)

        generate_content("Quantum Computing")
        call_kwargs = mock_client.messages.create.call_args[1]
        user_message = call_kwargs["messages"][0]["content"]
        self.assertIn("Quantum Computing", user_message)

    @patch("generate.anthropic.Anthropic")
    def test_skips_non_text_blocks(self, mock_anthropic_cls):
        thinking_block = MagicMock()
        thinking_block.type = "thinking"
        text_block = MagicMock()
        text_block.type = "text"
        text_block.text = json.dumps(SAMPLE_DATA)

        mock_client = MagicMock()
        mock_anthropic_cls.return_value = mock_client
        response = MagicMock()
        response.content = [thinking_block, text_block]
        mock_client.messages.create.return_value = response

        result = generate_content("Topic")
        self.assertEqual(result["title"], "Test Presentation")


# ---------------------------------------------------------------------------
# Output filename derivation (tested via main internals)
# ---------------------------------------------------------------------------

class TestFilenameDerivation(unittest.TestCase):

    def _derive(self, topic, fmt="pptx"):
        safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic)[:40].strip()
        return f"{safe}.{fmt}"

    def test_simple_topic(self):
        self.assertEqual(self._derive("Machine Learning"), "Machine Learning.pptx")

    def test_special_chars_replaced(self):
        name = self._derive("AI & ML: Future!")
        self.assertNotIn("&", name)
        self.assertNotIn(":", name)
        self.assertNotIn("!", name)

    def test_truncated_at_40_chars(self):
        long_topic = "A" * 50
        name = self._derive(long_topic)
        base = name.rsplit(".", 1)[0]
        self.assertLessEqual(len(base), 40)

    def test_html_extension(self):
        name = self._derive("My Topic", fmt="html")
        self.assertTrue(name.endswith(".html"))


# ---------------------------------------------------------------------------
# validate_topic tests
# ---------------------------------------------------------------------------

class TestValidateTopic(unittest.TestCase):

    def test_valid_topic_returned(self):
        self.assertEqual(validate_topic("Machine Learning"), "Machine Learning")

    def test_strips_whitespace(self):
        self.assertEqual(validate_topic("  AI  "), "AI")

    def test_empty_topic_exits(self):
        with self.assertRaises(SystemExit):
            validate_topic("")

    def test_whitespace_only_exits(self):
        with self.assertRaises(SystemExit):
            validate_topic("   ")

    def test_topic_at_max_length_passes(self):
        topic = "A" * TOPIC_MAX_LENGTH
        result = validate_topic(topic)
        self.assertEqual(result, topic)

    def test_topic_over_max_length_exits(self):
        with self.assertRaises(SystemExit):
            validate_topic("A" * (TOPIC_MAX_LENGTH + 1))


# ---------------------------------------------------------------------------
# validate_output_path tests
# ---------------------------------------------------------------------------

class TestValidateOutputPath(unittest.TestCase):

    def test_simple_filename_allowed(self):
        result = validate_output_path("output.pptx", "pptx")
        self.assertTrue(result.endswith("output.pptx"))

    def test_path_traversal_rejected(self):
        with self.assertRaises(SystemExit):
            validate_output_path("../../evil.pptx", "pptx")

    def test_absolute_path_outside_cwd_rejected(self):
        with self.assertRaises(SystemExit):
            validate_output_path("/tmp/evil.pptx", "pptx")

    def test_extension_added_if_missing(self):
        result = validate_output_path("myfile", "html")
        self.assertTrue(result.endswith(".html"))

    def test_subdirectory_within_cwd_allowed(self):
        import os
        cwd = os.getcwd()
        # A path that stays within cwd should be accepted
        result = validate_output_path("subdir/output.pptx", "pptx")
        self.assertTrue(result.startswith(cwd))


# ---------------------------------------------------------------------------
# _check_rate_limit tests
# ---------------------------------------------------------------------------

class TestCheckRateLimit(unittest.TestCase):

    def setUp(self):
        # Reset global state before each test
        import generate
        generate._last_api_call = 0.0

    def test_first_call_passes_immediately(self):
        start = time.monotonic()
        _check_rate_limit()
        self.assertLess(time.monotonic() - start, 1.0)

    def test_second_call_within_interval_sleeps(self):
        import generate
        generate._last_api_call = time.monotonic()  # simulate a very recent call
        with patch("generate.time.sleep") as mock_sleep:
            _check_rate_limit()
            mock_sleep.assert_called_once()
            wait = mock_sleep.call_args[0][0]
            self.assertGreater(wait, 0)

    def test_second_call_after_interval_passes_immediately(self):
        import generate
        generate._last_api_call = time.monotonic() - 20.0  # well past the interval
        with patch("generate.time.sleep") as mock_sleep:
            _check_rate_limit()
            mock_sleep.assert_not_called()

    def test_updates_last_api_call_timestamp(self):
        import generate
        before = time.monotonic()
        _check_rate_limit()
        self.assertGreaterEqual(generate._last_api_call, before)


if __name__ == "__main__":
    unittest.main()
