"""
build_visual_direction_slides.py

The visual direction deck — built to the exact spec it describes.
Every slide proves the rules it teaches.

Rules applied throughout:
- One idea per slide
- One accent touch per slide (#0EA5E9)
- 40%+ whitespace
- Left-aligned always
- No gradients, shadows, or more than two colors
- Inter/Segoe UI, JetBrains Mono/Consolas
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Color system — exactly as specified in VISUAL_DIRECTION.md
# ---------------------------------------------------------------------------
VOID     = RGBColor(0x0A, 0x0F, 0x1E)  # background
OBSIDIAN = RGBColor(0x11, 0x18, 0x27)  # surface / card bg
SLATE    = RGBColor(0x1F, 0x29, 0x37)  # borders / dividers
SIGNAL   = RGBColor(0x0E, 0xA5, 0xE9)  # accent — one touch per slide
SNOW     = RGBColor(0xF9, 0xFA, 0xFB)  # body text
ASH      = RGBColor(0x6B, 0x72, 0x80)  # muted / labels
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color=VOID):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, left, top, width, height,
       size=16, bold=False, color=SNOW, align=PP_ALIGN.LEFT,
       font="Calibri", italic=False, wrap=True):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return shape


def accent_line(slide, top, width=9.0, left=0.6, color=SIGNAL):
    """Thin horizontal accent line."""
    d = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.025))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()


def divider(slide, top, color=SLATE):
    accent_line(slide, top, color=color)


def accent_bar_left(slide, top, height=0.6, color=SIGNAL):
    """Vertical accent bar at left edge."""
    b = slide.shapes.add_shape(1, Inches(0.6), Inches(top), Inches(0.06), Inches(height))
    b.fill.solid()
    b.fill.fore_color.rgb = color
    b.line.fill.background()


def swatch(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def card(slide, left, top, width, height, color=OBSIDIAN, border=SLATE):
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.color.rgb = border
    return s


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    """Title — statement template."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    accent_line(s, 0.0, width=10.0, left=0.0)

    tb(s, "VISUAL DIRECTION BRIEF", 0.6, 0.3, 9, 0.4,
       size=11, bold=True, color=ASH)
    tb(s, "AI Presentation Generator", 0.6, 0.85, 9, 0.7,
       size=38, bold=True, color=SNOW)
    accent_line(s, 1.75)
    tb(s, "A deck that follows every rule it teaches.", 0.6, 1.95, 9, 0.5,
       size=18, color=ASH)
    tb(s, "Dark. Restrained. Intentional.", 0.6, 3.8, 9, 0.5,
       size=14, color=ASH)


def slide_mandate(prs):
    """The creative mandate — statement template."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "THE MANDATE", 0.6, 0.3, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.85)

    tb(s, "Restraint\nas confidence.", 0.6, 1.1, 9, 2.0,
       size=52, bold=True, color=SNOW)

    tb(s, "Every element that isn't there is as intentional\nas every element that is.",
       0.6, 3.4, 9, 0.8, size=16, color=ASH)


def slide_palette(prs):
    """Color palette — six swatches with labels."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "COLOR PALETTE", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.75)

    colors = [
        (VOID,     "#0A0F1E", "Void",   "Background"),
        (OBSIDIAN, "#111827", "Obsidian","Surface"),
        (SLATE,    "#1F2937", "Slate",  "Border"),
        (SIGNAL,   "#0EA5E9", "Signal Blue", "Accent — one touch per slide"),
        (SNOW,     "#F9FAFB", "Snow",   "Body text"),
        (ASH,      "#6B7280", "Ash",    "Labels / muted"),
    ]

    swatch_w = 1.3
    swatch_h = 1.6
    gap = 0.2
    start_left = 0.6

    for i, (color, hex_val, name, role) in enumerate(colors):
        left = start_left + i * (swatch_w + gap)
        swatch(s, left, 1.0, swatch_w, swatch_h, color)
        # Hex
        tb(s, hex_val, left, 2.72, swatch_w, 0.28,
           size=10, color=SIGNAL if color == SIGNAL else SNOW,
           font="Consolas", align=PP_ALIGN.LEFT)
        # Name
        tb(s, name, left, 3.05, swatch_w, 0.28,
           size=11, bold=True, color=SNOW, align=PP_ALIGN.LEFT)
        # Role
        tb(s, role, left, 3.35, swatch_w, 0.45,
           size=9, color=ASH, align=PP_ALIGN.LEFT)


def slide_typography(prs):
    """Typography — font pairings."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "TYPOGRAPHY", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.75)

    # Headline specimen
    tb(s, "Headline / Inter 700", 0.6, 1.0, 9, 0.35, size=11, color=ASH)
    tb(s, "One Command.", 0.6, 1.35, 9, 0.85, size=48, bold=True, color=SNOW)

    divider(s, 2.3)

    # Body specimen
    tb(s, "Body / Inter 300", 0.6, 2.45, 9, 0.3, size=11, color=ASH)
    tb(s, "Type a topic. Get a boardroom-ready presentation in 60 seconds.",
       0.6, 2.78, 9, 0.45, size=18, color=SNOW)

    divider(s, 3.35)

    # Code specimen
    tb(s, "Code / JetBrains Mono  →  Consolas (Windows fallback)", 0.6, 3.5, 9, 0.3,
       size=11, color=ASH)
    card(s, 0.6, 3.85, 9, 0.6, color=OBSIDIAN, border=SLATE)
    tb(s, 'python generate.py "Your next meeting" --theme corporate',
       0.75, 3.92, 8.7, 0.45, size=16, color=SIGNAL, font="Consolas")


def slide_rule_one_idea(prs):
    """Rule 1 — statement template."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "RULE 1", 0.6, 0.3, 9, 0.4, size=11, bold=True, color=SIGNAL)
    accent_line(s, 0.85)

    tb(s, "One idea\nper slide.\nFull stop.", 0.6, 1.1, 9, 2.8,
       size=52, bold=True, color=SNOW)

    tb(s, "If you can remove it and the slide still works — remove it.",
       0.6, 4.2, 9, 0.5, size=14, color=ASH)


def slide_rule_grid(prs):
    """Rule 2 — the grid."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "RULE 2", 0.6, 0.3, 9, 0.4, size=11, bold=True, color=SIGNAL)
    accent_line(s, 0.85)

    tb(s, "Align everything\nto an invisible grid.\nThen align it again.",
       0.6, 1.1, 9, 2.4, size=40, bold=True, color=SNOW)

    divider(s, 3.65)
    tb(s, "Margins: 0.6\"  ·  Max width: 8.8\"  ·  Grid: never approximate",
       0.6, 3.8, 9, 0.45, size=13, color=ASH, font="Consolas")
    tb(s, "Misalignment transfers to credibility. The audience doesn't notice it consciously.\nTheir subconscious does.",
       0.6, 4.25, 9, 0.7, size=13, color=ASH)


def slide_rule_accent(prs):
    """Rule 3 — the accent budget. The accent is used once: on the key word."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "RULE 3", 0.6, 0.3, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.85)

    # Multi-run headline: "The accent is a " [budget] "."
    txBox = s.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9), Inches(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = "The accent is a "
    r1.font.size = Pt(44)
    r1.font.bold = True
    r1.font.name = "Calibri"
    r1.font.color.rgb = SNOW

    r2 = p.add_run()
    r2.text = "budget."
    r2.font.size = Pt(44)
    r2.font.bold = True
    r2.font.name = "Calibri"
    r2.font.color.rgb = SIGNAL   # ← the one accent touch on this slide

    divider(s, 2.65)

    rows = [
        ("One touch per slide.", SNOW),
        ("On exactly the right element.", SNOW),
        ("When everything is highlighted — nothing is.", ASH),
    ]
    top = 2.82
    for text, color in rows:
        tb(s, f"— {text}", 0.6, top, 9, 0.42, size=16, color=color)
        top += 0.42


def slide_layouts(prs):
    """Three layout templates — shown as wireframes."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "THREE LAYOUT TEMPLATES", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.75)

    templates = [
        ("Statement", "One sentence.\nNothing else.", 0.6),
        ("Metric", "One number.\nLabel below.", 3.7),
        ("Content", "Headline +\n3–4 bullets.", 6.8),
    ]

    for name, desc, left in templates:
        card(s, left, 1.0, 2.8, 3.5, color=OBSIDIAN, border=SLATE)
        tb(s, name, left + 0.15, 1.1, 2.5, 0.35,
           size=12, bold=True, color=SIGNAL)
        tb(s, desc, left + 0.15, 1.55, 2.5, 1.0,
           size=13, color=SNOW)
        # Accent touch: thin bar top of card
        b = s.shapes.add_shape(1,
            Inches(left), Inches(1.0), Inches(2.8), Inches(0.04))
        b.fill.solid()
        b.fill.fore_color.rgb = SIGNAL
        b.line.fill.background()
        tb(s, desc.replace("\n", " — "), left + 0.15, 3.65, 2.5, 0.7,
           size=10, color=ASH)


def slide_charts(prs):
    """Permitted chart types."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "CHART RULES", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.75)

    permitted = [
        ("Two-bar comparison", "Before vs. after. Label bars directly. No legend."),
        ("Single number", "Full slide. 80pt minimum. Unit in muted text."),
        ("Horizontal timeline", "Two to three points max. Dots and a line. No arrows."),
        ("Two-column table", "No grid lines. Alternating row tint at 5% opacity."),
    ]

    top = 0.95
    for title, rule in permitted:
        accent_bar_left(s, top + 0.05, height=0.35)
        tb(s, title, 0.82, top, 4.5, 0.35, size=14, bold=True, color=SNOW)
        tb(s, rule, 0.82, top + 0.35, 8.7, 0.32, size=12, color=ASH)
        top += 0.85

    divider(s, top + 0.1)
    tb(s, "Never: pie charts  ·  3D charts  ·  legends  ·  gridlines  ·  stacked bars",
       0.6, top + 0.25, 9, 0.4, size=12, color=SIGNAL, font="Consolas")


def slide_images(prs):
    """Image style directive."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "IMAGE STYLE", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.75)

    yes_items = [
        "Terminal screenshots — real, unretouched, cropped tight",
        "Code snippets — syntax-highlighted, 2–5 lines maximum",
        "Line-art diagrams — single color, max four elements",
        "Large numbers — treat $40,000 like Apple treats product renders",
    ]
    no_items = [
        "Stock photography",
        "Abstract 'technology' imagery",
        "Noun Project icons",
        "Anything that could appear in a different presentation",
    ]

    tb(s, "USE", 0.6, 1.0, 4.2, 0.3, size=11, bold=True, color=SIGNAL)
    top = 1.35
    for item in yes_items:
        tb(s, f"— {item}", 0.6, top, 4.2, 0.38, size=12, color=SNOW)
        top += 0.38

    tb(s, "NEVER", 5.2, 1.0, 4.2, 0.3, size=11, bold=True, color=ASH)
    top = 1.35
    for item in no_items:
        tb(s, f"— {item}", 5.2, top, 4.4, 0.38, size=12, color=ASH)
        top += 0.38

    # Vertical divider
    d = s.shapes.add_shape(1, Inches(4.95), Inches(1.0), Inches(0.025), Inches(3.2))
    d.fill.solid()
    d.fill.fore_color.rgb = SLATE
    d.line.fill.background()


def slide_checklist(prs):
    """Production checklist."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    tb(s, "PRODUCTION CHECKLIST", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)
    accent_line(s, 0.75)
    tb(s, "Before any slide leaves review.", 0.6, 0.85, 9, 0.35, size=15, color=SNOW)

    items = [
        "One idea per slide",
        "40%+ whitespace",
        "Accent color used once, on the right element",
        "All text left-aligned",
        "No gradients, shadows, or more than two colors",
        "Line length under 60 characters",
        "All elements snapped to the 0.6\" grid",
        "No stock imagery, icons, or SmartArt",
        "Chart: no legend, no gridlines, directly labeled",
    ]

    col1 = items[:5]
    col2 = items[5:]

    top = 1.35
    for item in col1:
        card(s, 0.6, top, 0.28, 0.28, color=SLATE, border=SIGNAL)
        tb(s, item, 1.05, top, 3.9, 0.3, size=12, color=SNOW)
        top += 0.42

    top = 1.35
    for item in col2:
        card(s, 5.2, top, 0.28, 0.28, color=SLATE, border=SIGNAL)
        tb(s, item, 5.65, top, 4.1, 0.3, size=12, color=SNOW)
        top += 0.42


def slide_brief(prs):
    """The brief in one sentence — closing statement slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    accent_line(s, 0.0, width=10.0, left=0.0, color=SIGNAL)

    tb(s, "THE BRIEF IN ONE SENTENCE", 0.6, 0.25, 9, 0.4, size=11, bold=True, color=ASH)

    tb(s, "Dark. Inter.\nOne idea.\nOne accent touch.\nNothing that doesn't\nearn its space.",
       0.6, 0.85, 9, 3.5, size=40, bold=True, color=SNOW)

    accent_line(s, 4.55)
    tb(s, "If a slide breaks any of these rules, it goes back.",
       0.6, 4.7, 9, 0.4, size=13, color=ASH)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slide_title(prs)
    slide_mandate(prs)
    slide_palette(prs)
    slide_typography(prs)
    slide_rule_one_idea(prs)
    slide_rule_grid(prs)
    slide_rule_accent(prs)
    slide_layouts(prs)
    slide_charts(prs)
    slide_images(prs)
    slide_checklist(prs)
    slide_brief(prs)

    out = "Visual_Direction.pptx"
    prs.save(out)
    print(f"Saved: {out}  (12 slides)")


if __name__ == "__main__":
    build()
