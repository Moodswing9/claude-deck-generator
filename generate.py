"""
generate.py - AI-powered presentation generator.

Uses Claude to generate slide content from a topic, then outputs a styled
.pptx or .html file with your choice of color theme.

Usage:
    python generate.py "Your Topic" [--theme THEME] [--format FORMAT] [--output FILE]
    python generate.py "Your Topic" --images            # embed Unsplash photos
    python generate.py "Your Topic" --remix old.pptx   # remix an existing deck
    python generate.py "Your Topic" --slides 8         # control slide count
    python generate.py "Your Topic" --no-notes         # omit speaker notes

Themes:  dark (default), light, corporate, executive
Formats: pptx (default), html
"""

import argparse
import html
import io
import json
import os
import pathlib
import sys
import time
from datetime import datetime

import anthropic
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Themes
# ---------------------------------------------------------------------------

THEMES = {
    "dark": {
        "name": "Dark",
        # pptx colors
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0xE9, 0x45, 0x60),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtext": RGBColor(0xA0, 0xA0, 0xB0),
        "divider": RGBColor(0x0F, 0x34, 0x60),
        # html colors
        "background": "#1a1a2e",
        "slide_bg": "#16213e",
        "primary": "#e94560",
        "secondary": "#0f3460",
        "html_text": "#eaeaea",
        "muted": "#a0a0b0",
        "border": "#0f3460",
        "code_bg": "#0d0d1a",
        "font_family": "'Segoe UI', Tahoma, Geneva, Verdana, sans-serif",
    },
    "light": {
        "name": "Light",
        "bg": RGBColor(0xF5, 0xF5, 0xF5),
        "accent": RGBColor(0x25, 0x63, 0xEB),
        "text": RGBColor(0x11, 0x18, 0x27),
        "subtext": RGBColor(0x6B, 0x72, 0x80),
        "divider": RGBColor(0xD1, 0xD5, 0xDB),
        "background": "#f5f5f5",
        "slide_bg": "#ffffff",
        "primary": "#2563eb",
        "secondary": "#e5e7eb",
        "html_text": "#111827",
        "muted": "#6b7280",
        "border": "#d1d5db",
        "code_bg": "#f3f4f6",
        "font_family": "'Segoe UI', Tahoma, Geneva, Verdana, sans-serif",
    },
    "corporate": {
        "name": "Corporate",
        "bg": RGBColor(0x1E, 0x29, 0x3B),
        "accent": RGBColor(0x0E, 0xA5, 0xE9),
        "text": RGBColor(0xF1, 0xF5, 0xF9),
        "subtext": RGBColor(0x94, 0xA3, 0xB8),
        "divider": RGBColor(0x33, 0x41, 0x55),
        "background": "#1e293b",
        "slide_bg": "#0f172a",
        "primary": "#0ea5e9",
        "secondary": "#1e40af",
        "html_text": "#f1f5f9",
        "muted": "#94a3b8",
        "border": "#334155",
        "code_bg": "#0a0f1e",
        "font_family": "'Inter', 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif",
    },
    "executive": {
        "name": "Executive",
        # pptx colors — warm off-white, deep navy, gold accent
        "bg": RGBColor(0xF5, 0xF4, 0xF0),
        "accent": RGBColor(0xB8, 0x86, 0x00),
        "text": RGBColor(0x0D, 0x1B, 0x2A),
        "subtext": RGBColor(0x55, 0x65, 0x7A),
        "divider": RGBColor(0xE2, 0xDE, 0xD5),
        # html colors
        "background": "#eeece6",
        "slide_bg": "#f5f4f0",
        "primary": "#b88600",
        "secondary": "#ece9e1",
        "html_text": "#0d1b2a",
        "muted": "#55657a",
        "border": "#e2ded5",
        "code_bg": "#f0ede8",
        "font_family": "'Georgia', 'Times New Roman', serif",
    },
}

DEFAULT_THEME = "dark"

TOPIC_MAX_LENGTH = 200
SLIDES_MIN = 4
SLIDES_MAX = 20
SLIDES_DEFAULT = 12

# ---------------------------------------------------------------------------
# Input validation
# ---------------------------------------------------------------------------

def validate_topic(topic: str) -> str:
    """Validate and sanitize the presentation topic."""
    topic = topic.strip()
    if not topic:
        raise ValueError("Topic is required.")
    if len(topic) > TOPIC_MAX_LENGTH:
        raise ValueError(
            f"Topic must be {TOPIC_MAX_LENGTH} characters or fewer (got {len(topic)})."
        )
    return topic


def validate_output_path(path: str, fmt: str) -> str:
    """Reject paths that escape the current working directory."""
    cwd = pathlib.Path.cwd().resolve()
    resolved = (cwd / path).resolve()
    if not str(resolved).startswith(str(cwd)):
        raise ValueError(f"Output path '{path}' is outside the current directory.")
    if not resolved.suffix:
        resolved = resolved.with_suffix(f".{fmt}")
    return str(resolved)


# ---------------------------------------------------------------------------
# Rate limiter
# ---------------------------------------------------------------------------

_last_api_call: float = 0.0
_MIN_INTERVAL = 10.0  # minimum seconds between API calls


def _check_rate_limit():
    global _last_api_call
    elapsed = time.monotonic() - _last_api_call
    if _last_api_call and elapsed < _MIN_INTERVAL:
        wait = _MIN_INTERVAL - elapsed
        print(f"Rate limit: waiting {wait:.1f}s before next API call...")
        time.sleep(wait)
    _last_api_call = time.monotonic()


# ---------------------------------------------------------------------------
# Unsplash image fetcher
# ---------------------------------------------------------------------------

UNSPLASH_ACCESS_KEY = os.environ.get("UNSPLASH_ACCESS_KEY", "")
_UNSPLASH_API = "https://api.unsplash.com/search/photos"


def fetch_slide_images(slides: list) -> dict:
    """Return {slide_title: image_url} for each slide using Unsplash search."""
    if not UNSPLASH_ACCESS_KEY:
        print("Warning: UNSPLASH_ACCESS_KEY not set — skipping images.")
        return {}

    images = {}
    headers = {"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"}
    for slide in slides:
        title = slide.get("title", "")
        if not title or slide.get("type") == "section":
            continue
        try:
            resp = requests.get(
                _UNSPLASH_API,
                headers=headers,
                params={"query": title, "per_page": 1, "orientation": "landscape"},
                timeout=6,
            )
            if resp.status_code == 200:
                results = resp.json().get("results", [])
                if results:
                    images[title] = results[0]["urls"]["small"]
        except Exception:
            pass
    return images


def _download_image(url: str) -> io.BytesIO | None:
    """Download an image URL into a BytesIO buffer."""
    try:
        resp = requests.get(url, timeout=10)
        if resp.status_code == 200:
            return io.BytesIO(resp.content)
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# MarkItDown — ingest existing PPTX as Markdown context
# ---------------------------------------------------------------------------

def ingest_pptx(source: str) -> str:
    """
    Convert an existing PPTX file to Markdown using MarkItDown.

    Parameters:
    - source: Path to the .pptx file to ingest.

    Returns:
    - Markdown string extracted from the presentation.
    """
    try:
        from markitdown import MarkItDown
    except ImportError:
        raise ImportError(
            "markitdown is required for the remix feature. "
            "Install it with: pip install 'markitdown[pptx]'"
        )
    result = MarkItDown().convert(source)
    return result.markdown


# ---------------------------------------------------------------------------
# Claude API — generate slide content
# ---------------------------------------------------------------------------

SLIDE_SCHEMA = {
    "type": "object",
    "properties": {
        "title": {"type": "string"},
        "subtitle": {"type": "string"},
        "slides": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "type": {
                        "type": "string",
                        "enum": ["content", "section", "quote", "stat"],
                    },
                    "title": {"type": "string"},
                    "bullets": {"type": "array", "items": {"type": "string"}},
                    "quote": {"type": "string"},
                    "attribution": {"type": "string"},
                    "stat": {"type": "string"},
                    "stat_label": {"type": "string"},
                    "notes": {"type": "string"},
                },
                "required": ["type", "title", "bullets", "notes"],
                "additionalProperties": False,
            },
        },
    },
    "required": ["title", "subtitle", "slides"],
    "additionalProperties": False,
}


def _nvidia_client():
    """Return an OpenAI-compatible client pointed at NVIDIA NIM."""
    try:
        from openai import OpenAI as _OpenAI
    except ImportError:
        raise ImportError("openai>=1.0 is required for NVIDIA provider. Install it: pip install openai>=1.0")
    return _OpenAI(
        base_url="https://integrate.api.nvidia.com/v1",
        api_key=os.environ["NVIDIA_API_KEY"],
    )


def analyze_slide_image(image_path: str) -> str:
    """Use Phi-4 Multimodal to describe a slide image and extract content."""
    import base64
    client = _nvidia_client()
    with open(image_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    ext = pathlib.Path(image_path).suffix.lstrip(".").lower()
    mime = "image/png" if ext == "png" else "image/jpeg"
    resp = client.chat.completions.create(
        model="microsoft/phi-4-multimodal-instruct",
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                {"type": "text", "text": (
                    "Describe this presentation slide in detail. "
                    "Extract the title, all visible text, and any data or statistics shown."
                )},
            ],
        }],
        max_tokens=512,
    )
    return resp.choices[0].message.content or ""


def extract_chart_data(image_path: str) -> str:
    """Use DePlot to extract data table from a chart image."""
    import base64
    client = _nvidia_client()
    with open(image_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode()
    ext = pathlib.Path(image_path).suffix.lstrip(".").lower()
    mime = "image/png" if ext == "png" else "image/jpeg"
    resp = client.chat.completions.create(
        model="google/deplot",
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                {"type": "text", "text": "Generate underlying data table of the figure below:"},
            ],
        }],
        max_tokens=512,
    )
    return resp.choices[0].message.content or ""


def generate_content_nvidia(
    topic: str,
    *,
    reference_markdown: str = "",
    slide_count: int = SLIDES_DEFAULT,
) -> dict:
    """Generate slide content using Writer Palmyra-Creative-122B via NVIDIA NIM.

    Returns the same dict structure as generate_content() so downstream
    build_pptx() / build_html() work unchanged.
    """
    _check_rate_limit()
    slide_count = max(SLIDES_MIN, min(SLIDES_MAX, slide_count))
    client = _nvidia_client()

    if reference_markdown:
        print(f"[NVIDIA] Remixing deck into {slide_count} slides about: {topic}")
        user_message = (
            f"Here is an existing presentation for reference:\n\n"
            f"<reference_deck>\n{reference_markdown}\n</reference_deck>\n\n"
            f"Using the above as source material, create an improved, professional "
            f"{slide_count}-slide presentation about: {topic}\n\n"
            "Preserve strong points, sharpen language, improve narrative structure. "
            "Structure: opening hook, agenda, 2-3 major sections, closing."
        )
    else:
        print(f"[NVIDIA] Generating {slide_count}-slide presentation about: {topic}")
        user_message = (
            f"Create a professional {slide_count}-slide presentation about: {topic}\n\n"
            "Structure: opening hook (stat or quote), agenda, 2-3 major sections each "
            "preceded by a section slide, supporting content slides, and a strong closing."
        )

    system_prompt = (
        "You are a McKinsey-level presentation strategist. Create executive-quality "
        "slide decks with clear narrative flow, precise language, and varied layouts.\n\n"
        "Return ONLY a valid JSON object matching this exact schema — no markdown fences, no prose:\n"
        '{"title": "...", "subtitle": "...", "slides": [{"type": "content|section|quote|stat", '
        '"title": "...", "bullets": [], "notes": "...", '
        '"quote": "...", "attribution": "...", "stat": "...", "stat_label": "..."}]}\n\n'
        "Slide type rules:\n"
        "- 'content': 3-5 bullet points each under 12 words\n"
        "- 'section': transition slide, bullets=[]\n"
        "- 'quote': set bullets=[], fill quote and attribution\n"
        "- 'stat': set bullets=[], fill stat and stat_label\n"
        "Every 2-3 content slides, insert a section/quote/stat for visual rhythm."
    )

    resp = client.chat.completions.create(
        model="writer/palmyra-creative-122b",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_message},
        ],
        temperature=0.7,
        max_tokens=4096,
    )
    raw = resp.choices[0].message.content or "{}"
    # Strip markdown fences if the model wraps anyway
    raw = raw.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw)


def generate_content(
    topic: str,
    *,
    reference_markdown: str = "",
    slide_count: int = SLIDES_DEFAULT,
) -> dict:
    """
    Call Claude to generate structured slide content.

    Parameters:
    - topic: The presentation subject.
    - reference_markdown: Optional Markdown from an existing PPTX (via ingest_pptx).
                          When provided, Claude uses it as source material to remix.
    - slide_count: Target number of slides (clamped to SLIDES_MIN–SLIDES_MAX).
    """
    _check_rate_limit()
    slide_count = max(SLIDES_MIN, min(SLIDES_MAX, slide_count))
    client = anthropic.Anthropic()

    if reference_markdown:
        print(f"Remixing existing deck into {slide_count} slides about: {topic}")
        user_message = (
            f"Here is an existing presentation for reference:\n\n"
            f"<reference_deck>\n{reference_markdown}\n</reference_deck>\n\n"
            f"Using the above as source material, create an improved, professional "
            f"{slide_count}-slide presentation about: {topic}\n\n"
            "Preserve strong points from the reference, sharpen the language, improve "
            "narrative structure, and ensure every slide earns its place. "
            "Structure: opening hook (stat or quote), agenda, 2-3 major sections each "
            "preceded by a section slide, supporting content slides, and a strong closing."
        )
    else:
        print(f"Generating {slide_count}-slide presentation about: {topic}")
        user_message = (
            f"Create a professional {slide_count}-slide presentation about: {topic}\n\n"
            "Structure: opening hook (stat or quote), agenda, 2-3 major sections each "
            "preceded by a section slide, supporting content slides, and a strong closing."
        )

    response = client.messages.create(
        model="claude-opus-4-6",
        max_tokens=6000,
        thinking={"type": "adaptive"},
        system=(
            "You are a McKinsey-level presentation strategist. Create executive-quality "
            "slide decks with clear narrative flow, precise language, and varied layouts.\n\n"
            "Slide type guidelines:\n"
            "- 'content': 3-5 crisp bullet points, each under 12 words\n"
            "- 'section': transition slide between major sections, no bullets needed, "
            "  set bullets to []\n"
            "- 'quote': a compelling quote relevant to the topic, set bullets to [], "
            "  populate quote and attribution fields\n"
            "- 'stat': one striking statistic or number (e.g. '$4.2T', '73%', '10x'), "
            "  set bullets to [], populate stat and stat_label fields\n\n"
            "Use a mix of types. Every 2-3 content slides, insert a section, quote, or stat "
            "slide to maintain visual rhythm. Avoid filler — every slide must earn its place."
        ),
        messages=[{"role": "user", "content": user_message}],
        output_config={
            "format": {"type": "json_schema", "schema": SLIDE_SCHEMA}
        },
    )
    text = next(b.text for b in response.content if b.type == "text")
    return json.loads(text)


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

def _pptx_set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _pptx_textbox(slide, text, left, top, width, height,
                  size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return tb


def _pptx_title_slide(prs, title, subtitle, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_set_bg(slide, theme["bg"])
    # Left accent bar
    bar = slide.shapes.add_shape(1, 0, Inches(1.5), Inches(0.5), Inches(2.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()
    # Bottom stripe
    stripe = slide.shapes.add_shape(1, 0, Inches(5.1), Inches(10), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = theme["accent"]
    stripe.line.fill.background()
    _pptx_textbox(slide, title, 0.7, 1.5, 8.6, 1.8, size=44, bold=True,
                  color=theme["text"], align=PP_ALIGN.LEFT)
    _pptx_textbox(slide, subtitle or "AI-Generated Presentation", 0.7, 3.4, 8.6, 0.7,
                  size=20, color=theme["subtext"], align=PP_ALIGN.LEFT)


def _pptx_content_slide(prs, title, bullets, notes, theme, image_stream=None,
                         include_notes: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_set_bg(slide, theme["bg"])

    # With image: text on left 55%, image on right 40%
    text_width = 5.5 if image_stream else 8.8

    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(0.06), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()

    _pptx_textbox(slide, title, 0.7, 0.28, text_width, 0.9, size=28, bold=True, color=theme["text"])

    div = slide.shapes.add_shape(1, Inches(0.5), Inches(1.12), Inches(text_width + 0.2), Inches(0.02))
    div.fill.solid()
    div.fill.fore_color.rgb = theme["divider"]
    div.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(text_width), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size = Pt(19)
        run.font.name = "Calibri"
        run.font.color.rgb = theme["text"]

    if image_stream:
        slide.shapes.add_picture(image_stream, Inches(6.4), Inches(1.2), Inches(3.2), Inches(3.8))

    if include_notes and notes:
        slide.notes_slide.notes_text_frame.text = notes


def _pptx_section_slide(prs, title, notes, theme, include_notes: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_set_bg(slide, theme["accent"])
    # Full-bleed accent background, white title
    _pptx_textbox(slide, title, 0.8, 1.8, 8.4, 2.0, size=38, bold=True,
                  color=theme["bg"], align=PP_ALIGN.LEFT)
    if include_notes and notes:
        slide.notes_slide.notes_text_frame.text = notes


def _pptx_quote_slide(prs, title, quote, attribution, notes, theme, include_notes: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_set_bg(slide, theme["bg"])

    # Large quotation mark
    _pptx_textbox(slide, "\u201c", 0.4, 0.1, 1.5, 1.2, size=72, bold=True,
                  color=theme["accent"], align=PP_ALIGN.LEFT)

    quote_text = quote or title
    _pptx_textbox(slide, quote_text, 0.7, 1.0, 8.6, 2.8, size=22,
                  color=theme["text"], align=PP_ALIGN.LEFT)

    if attribution:
        _pptx_textbox(slide, f"\u2014 {attribution}", 0.7, 4.0, 8.6, 0.5,
                      size=16, color=theme["subtext"], align=PP_ALIGN.LEFT)

    # Bottom accent line
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(4.7), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = theme["accent"]
    line.line.fill.background()

    if include_notes and notes:
        slide.notes_slide.notes_text_frame.text = notes


def _pptx_stat_slide(prs, title, stat, stat_label, notes, theme, include_notes: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_set_bg(slide, theme["bg"])

    _pptx_textbox(slide, title, 0.7, 0.25, 8.6, 0.7, size=22, bold=True,
                  color=theme["subtext"], align=PP_ALIGN.CENTER)

    div = slide.shapes.add_shape(1, Inches(4.0), Inches(1.05), Inches(2.0), Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = theme["accent"]
    div.line.fill.background()

    _pptx_textbox(slide, stat or "", 0.5, 1.2, 9.0, 2.4, size=80, bold=True,
                  color=theme["accent"], align=PP_ALIGN.CENTER)

    _pptx_textbox(slide, stat_label or "", 0.5, 3.7, 9.0, 0.8, size=20,
                  color=theme["text"], align=PP_ALIGN.CENTER)

    if include_notes and notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_pptx(data: dict, theme: dict, output_path: str, images: dict = None,
               include_notes: bool = True):
    images = images or {}
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    _pptx_title_slide(prs, data["title"], data.get("subtitle", ""), theme)
    for s in data["slides"]:
        stype = s.get("type", "content")
        notes = s.get("notes", "")
        if stype == "section":
            _pptx_section_slide(prs, s["title"], notes, theme, include_notes=include_notes)
        elif stype == "quote":
            _pptx_quote_slide(prs, s["title"], s.get("quote", ""),
                              s.get("attribution", ""), notes, theme, include_notes=include_notes)
        elif stype == "stat":
            _pptx_stat_slide(prs, s["title"], s.get("stat", ""),
                             s.get("stat_label", ""), notes, theme, include_notes=include_notes)
        else:
            image_stream = _download_image(images[s["title"]]) if s["title"] in images else None
            _pptx_content_slide(prs, s["title"], s["bullets"], notes, theme, image_stream,
                                include_notes=include_notes)
    prs.save(output_path)


# ---------------------------------------------------------------------------
# HTML builder
# ---------------------------------------------------------------------------

def _css(theme: dict) -> str:
    t = theme
    return f"""
        :root {{
            --bg: {t['background']};
            --slide-bg: {t['slide_bg']};
            --primary: {t['primary']};
            --secondary: {t['secondary']};
            --text: {t['html_text']};
            --muted: {t['muted']};
            --border: {t['border']};
            --code-bg: {t['code_bg']};
        }}
        * {{ box-sizing: border-box; margin: 0; padding: 0; }}
        body {{
            background: var(--bg); color: var(--text);
            font-family: {t['font_family']};
            min-height: 100vh; display: flex; flex-direction: column;
            align-items: center; padding: 2rem 1rem;
        }}
        h1, h2 {{ color: var(--primary); }}
        .deck-title {{
            text-align: center; margin-bottom: 2.5rem; padding-bottom: 1.2rem;
            border-bottom: 2px solid var(--border); width: 100%; max-width: 900px;
        }}
        .deck-title h1 {{ font-size: 2.2rem; letter-spacing: -0.02em; }}
        .deck-title .subtitle {{ color: var(--muted); font-size: 1.1rem; margin-top: 0.5rem; }}
        .deck-title .meta {{ color: var(--muted); font-size: 0.82rem; margin-top: 0.6rem; }}
        .theme-badge {{
            display: inline-block; background: var(--secondary); color: var(--primary);
            border-radius: 20px; padding: 0.2rem 0.9rem; font-size: 0.78rem;
            font-weight: 600; letter-spacing: 0.05em; margin-top: 0.5rem;
        }}
        /* Base slide */
        .slide {{
            background: var(--slide-bg); border: 1px solid var(--border);
            border-radius: 12px; width: 100%; max-width: 900px; min-height: 300px;
            margin-bottom: 1.5rem; padding: 2.8rem 3.5rem;
            display: flex; flex-direction: column; justify-content: center;
            position: relative; box-shadow: 0 2px 16px rgba(0,0,0,0.08);
            border-left: 4px solid transparent;
        }}
        .slide-number {{ position: absolute; top: 1rem; right: 1.5rem; font-size: 0.72rem; color: var(--muted); letter-spacing: 0.05em; }}
        /* Title slide */
        .slide-type-title {{
            border-left: 4px solid var(--primary);
            min-height: 340px;
        }}
        .slide-type-title h2 {{ font-size: 2.6rem; letter-spacing: -0.02em; margin-bottom: 0.8rem; line-height: 1.2; }}
        .slide-type-title .subtitle {{ color: var(--muted); font-size: 1.15rem; margin-top: 0.3rem; }}
        /* Section slide */
        .slide-type-section {{
            background: var(--primary); border-color: var(--primary);
            border-left: 4px solid rgba(0,0,0,0.15);
        }}
        .slide-type-section h2 {{ color: var(--slide-bg); font-size: 2.2rem; letter-spacing: -0.01em; }}
        .slide-type-section .slide-number {{ color: rgba(255,255,255,0.5); }}
        /* Content slide */
        .slide-type-content h2 {{
            font-size: 1.7rem; margin-bottom: 1.4rem; padding-bottom: 0.6rem;
            border-bottom: 1px solid var(--border);
        }}
        ul {{ list-style: none; padding: 0; }}
        ul li {{
            padding: 0.5rem 0 0.5rem 1.8rem; position: relative;
            font-size: 1.08rem; line-height: 1.65;
        }}
        ul li::before {{ content: '▸'; position: absolute; left: 0; color: var(--primary); font-size: 0.85em; top: 0.6rem; }}
        /* Quote slide */
        .slide-type-quote {{ border-left: 4px solid var(--primary); }}
        .slide-type-quote .quote-mark {{ font-size: 5rem; color: var(--primary); line-height: 0.8; margin-bottom: 0.5rem; opacity: 0.6; }}
        .slide-type-quote .quote-text {{ font-size: 1.35rem; line-height: 1.7; font-style: italic; margin-bottom: 1.2rem; }}
        .slide-type-quote .attribution {{ color: var(--muted); font-size: 0.95rem; border-top: 1px solid var(--border); padding-top: 0.8rem; }}
        /* Stat slide */
        .slide-type-stat {{ text-align: center; }}
        .slide-type-stat h2 {{ font-size: 1.2rem; color: var(--muted); margin-bottom: 1rem; text-transform: uppercase; letter-spacing: 0.1em; font-weight: 500; }}
        .slide-type-stat .stat-number {{ font-size: 5.5rem; font-weight: 700; color: var(--primary); line-height: 1; margin-bottom: 0.6rem; letter-spacing: -0.03em; }}
        .slide-type-stat .stat-label {{ font-size: 1.2rem; color: var(--text); }}
        /* Closing slide */
        .slide-type-closing {{
            text-align: center; border-top: 4px solid var(--primary);
            border-left: none; border-radius: 12px;
        }}
        .slide-type-closing h2 {{ font-size: 2.6rem; margin-bottom: 0.75rem; }}
        /* Image layout */
        .slide-has-image .slide-body {{ display: flex; gap: 2rem; align-items: center; }}
        .slide-has-image .slide-text {{ flex: 1; }}
        .slide-has-image .slide-image {{ flex: 0 0 38%; }}
        .slide-has-image .slide-image img {{ width: 100%; border-radius: 8px; object-fit: cover; max-height: 220px; }}
    """


def _slide_html(slide: dict, index: int, total: int) -> str:
    title = html.escape(slide.get("title", ""))
    bullets = slide.get("bullets", [])
    stype = slide.get("type", "content")
    num = f'<span class="slide-number">{index} / {total}</span>'
    items = "".join(f"<li>{html.escape(b)}</li>" for b in bullets)

    if index == 1:  # title slide
        subtitle = html.escape(slide.get("subtitle", ""))
        sub = f'<p class="subtitle">{subtitle}</p>' if subtitle else ""
        return (f'<section class="slide slide-type-title">{num}'
                f'<h2>{title}</h2>{sub}</section>')

    if index == total:  # closing slide
        return (f'<section class="slide slide-type-closing">{num}'
                f'<h2>{title}</h2></section>')

    if stype == "section":
        return (f'<section class="slide slide-type-section">{num}'
                f'<h2>{title}</h2></section>')

    if stype == "quote":
        quote = html.escape(slide.get("quote", slide.get("title", "")))
        attr = html.escape(slide.get("attribution", ""))
        attr_html = f'<p class="attribution">&mdash; {attr}</p>' if attr else ""
        return (f'<section class="slide slide-type-quote">{num}'
                f'<div class="quote-mark">&ldquo;</div>'
                f'<p class="quote-text">{quote}</p>{attr_html}</section>')

    if stype == "stat":
        stat = html.escape(slide.get("stat", ""))
        label = html.escape(slide.get("stat_label", ""))
        return (f'<section class="slide slide-type-stat">{num}'
                f'<h2>{title}</h2>'
                f'<div class="stat-number">{stat}</div>'
                f'<div class="stat-label">{label}</div></section>')

    image_url = slide.get("image_url", "")
    if image_url:
        safe_url = html.escape(image_url, quote=True)
        return (f'<section class="slide slide-type-content slide-has-image">{num}'
                f'<div class="slide-body">'
                f'<div class="slide-text"><h2>{title}</h2><ul>{items}</ul></div>'
                f'<div class="slide-image"><img src="{safe_url}" alt="{title}"/></div>'
                f'</div></section>')
    return (f'<section class="slide slide-type-content">{num}'
            f'<h2>{title}</h2><ul>{items}</ul></section>')


def build_html(data: dict, theme: dict, output_path: str, images: dict = None,
               include_notes: bool = True):
    images = images or {}
    content_slides = [
        {**s, "image_url": images.get(s["title"], "")} for s in data["slides"]
    ]
    title_slide = {"title": data["title"], "subtitle": data.get("subtitle", ""), "type": "title"}
    slides = [title_slide] + content_slides
    total = len(slides)
    slides_html = "".join(_slide_html(s, i + 1, total) for i, s in enumerate(slides))
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M")
    deck_title = html.escape(data["title"])
    deck_subtitle = html.escape(data.get("subtitle", ""))
    output_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8"/>
    <meta name="viewport" content="width=device-width, initial-scale=1.0"/>
    <title>{deck_title}</title>
    <style>{_css(theme)}</style>
</head>
<body>
    <div class="deck-title">
        <h1>{deck_title}</h1>
        {f'<p class="subtitle">{deck_subtitle}</p>' if deck_subtitle else ""}
        <div class="meta">Generated on {generated_at} &bull;
            <span class="theme-badge">{theme['name']} theme</span>
        </div>
    </div>
    {slides_html}
</body>
</html>"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(output_html)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Generate a presentation from a topic using Claude AI."
    )
    parser.add_argument("topic", nargs="?", help="Presentation topic")
    parser.add_argument("--theme", choices=list(THEMES.keys()), default=DEFAULT_THEME,
                        help=f"Color theme (default: {DEFAULT_THEME})")
    parser.add_argument("--format", choices=["pptx", "html"], default="pptx",
                        help="Output format (default: pptx)")
    parser.add_argument("--output", default=None,
                        help="Output file path (default: auto-named from topic)")
    parser.add_argument("--list-themes", action="store_true",
                        help="List available themes and exit")
    parser.add_argument("--images", action="store_true",
                        help="Embed Unsplash photos (requires UNSPLASH_ACCESS_KEY)")
    parser.add_argument("--remix", metavar="FILE", default=None,
                        help="Path to an existing .pptx — MarkItDown extracts it, Claude rebuilds it")
    parser.add_argument("--slides", type=int, default=SLIDES_DEFAULT, metavar="N",
                        help=f"Target slide count ({SLIDES_MIN}–{SLIDES_MAX}, default: {SLIDES_DEFAULT})")
    parser.add_argument("--no-notes", action="store_true",
                        help="Omit speaker notes from the output")
    parser.add_argument("--provider", choices=["anthropic", "nvidia"], default="anthropic",
                        help="AI provider for content generation (default: anthropic)")
    args = parser.parse_args()

    if args.list_themes:
        for key, t in THEMES.items():
            marker = " (default)" if key == DEFAULT_THEME else ""
            print(f"  {key:<12} {t['name']}{marker}")
        return

    try:
        topic = validate_topic(args.topic or input("Enter presentation topic: ").strip())
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)

    reference_markdown = ""
    if args.remix:
        if not args.remix.lower().endswith(".pptx"):
            print("Error: --remix requires a .pptx file.")
            sys.exit(1)
        print(f"Ingesting reference deck: {args.remix}")
        try:
            reference_markdown = ingest_pptx(args.remix)
        except (FileNotFoundError, ImportError) as e:
            print(f"Error: {e}")
            sys.exit(1)

    theme = THEMES[args.theme]
    if args.provider == "nvidia":
        if not os.environ.get("NVIDIA_API_KEY"):
            print("Error: NVIDIA_API_KEY environment variable is required for --provider nvidia")
            sys.exit(1)
        data = generate_content_nvidia(
            topic,
            reference_markdown=reference_markdown,
            slide_count=args.slides,
        )
    else:
        data = generate_content(
            topic,
            reference_markdown=reference_markdown,
            slide_count=args.slides,
        )

    images = {}
    if args.images:
        print("Fetching Unsplash images...")
        images = fetch_slide_images(data["slides"])
        print(f"Found images for {len(images)} slides.")

    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic)[:40].strip()
    raw_output = args.output or f"{safe}.{args.format}"
    try:
        output = validate_output_path(raw_output, args.format)
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)

    include_notes = not args.no_notes
    if args.format == "pptx":
        build_pptx(data, theme, output, images, include_notes=include_notes)
    else:
        build_html(data, theme, output, images, include_notes=include_notes)

    notes_note = " (no speaker notes)" if not include_notes else ""
    remix_note = f" [remixed from {args.remix}]" if args.remix else ""
    provider_note = f" via {args.provider}"
    print(f"Saved: {output}  ({len(data['slides'])} slides, {args.theme} theme{notes_note}{remix_note}{provider_note})")


if __name__ == "__main__":
    main()
