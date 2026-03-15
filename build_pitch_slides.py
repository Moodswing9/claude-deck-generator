"""
build_pitch_slides.py — generates AI_Presentation_Generator_Pitch.pptx

Content mirrors PRESENTATION_BLUEPRINT.md slide-for-slide.
Theme: corporate dark (navy/slate/sky blue).
No unsourced statistics. No fabricated compliance claims. No roadmap.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colors
# ---------------------------------------------------------------------------
BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
SLIDE   = RGBColor(0x1E, 0x29, 0x3B)   # slate
ACCENT  = RGBColor(0x0E, 0xA5, 0xE9)   # sky blue
TEXT    = RGBColor(0xF1, 0xF5, 0xF9)   # near white
MUTED   = RGBColor(0x94, 0xA3, 0xB8)   # slate gray
DIVIDER = RGBColor(0x33, 0x41, 0x55)   # dark border
RED     = RGBColor(0xEF, 0x44, 0x44)   # hook / urgency
GREEN   = RGBColor(0x22, 0xC5, 0x5E)   # positive signal

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def accent_bar(slide, color=ACCENT):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def divider(slide, top, color=DIVIDER):
    d = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(9), Inches(0.03))
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()


def textbox(slide, text, left, top, width, height,
            size=15, bold=False, color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return tb


def bullets(slide, items, left, top, width, height, size=15, color=TEXT, dot_color=ACCENT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5) if i > 0 else Pt(0)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = f"\u25b8  {item}"
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = color
    return tb


def code_block(slide, text, left, top, width, height):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x0A, 0x0F, 0x1E)
    box.line.color.rgb = DIVIDER
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.12),
                                   Inches(width - 0.4), Inches(height - 0.24))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(17)
    run.font.name = "Consolas"
    run.font.bold = True
    run.font.color.rgb = ACCENT
    return tb

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def s1_hook(prs):
    """Slide 1 — The Hook"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide, RED)

    textbox(slide,
            "You spent 4 hours on your last presentation.",
            0.6, 1.6, 8.8, 2.4,
            size=40, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

    textbox(slide, "1 / 11", 8.8, 5.2, 1, 0.3,
            size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s2_real_cost(prs):
    """Slide 2 — The Real Cost"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    textbox(slide, "You know the number.\nYou just haven't said it out loud.",
            0.5, 0.15, 9, 0.9, size=22, bold=True, color=TEXT)
    divider(slide, 1.18)

    items = [
        "Ask the room: how long did your last important presentation take to build?",
        "Three hours? Four? Six, if it was a board deck?",
        "Now multiply that by every person in this building who builds decks.",
        "Multiply by 52 weeks.",
        "That number doesn't appear on your P&L.  It should.",
    ]
    bullets(slide, items, 0.5, 1.3, 9, 3.5, size=16)
    textbox(slide, "2 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s3_status_quo(prs):
    """Slide 3 — The Status Quo Is a Choice"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    textbox(slide, "You already know how to fix this.\nYou just haven't done it.",
            0.5, 0.15, 9, 0.9, size=22, bold=True, color=TEXT)
    divider(slide, 1.18)

    items = [
        "Templates don't help — you still write every word",
        "Canva and Slides are faster tools for the same slow process",
        "AI assistants give you a blank page with a chatbot attached",
    ]
    bullets(slide, items, 0.5, 1.3, 9, 2.2, size=17)
    textbox(slide, "3 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s4_premise(prs):
    """Slide 4 — The Premise"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide, ACCENT)

    textbox(slide, "What if the presentation wrote itself?",
            0.6, 1.5, 8.8, 1.2,
            size=36, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    textbox(slide,
            "Not assisted. Not templated. Written — from a single sentence you already know.",
            0.6, 2.9, 8.8, 0.7,
            size=17, color=MUTED, align=PP_ALIGN.CENTER)
    textbox(slide, "4 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s5_demo(prs):
    """Slide 5 — The Demonstration"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide)

    textbox(slide, "Watch.", 0.6, 0.9, 8.8, 1.2,
            size=60, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    code_block(slide, 'python generate.py "Your Topic" --theme corporate',
               0.8, 2.5, 8.4, 0.65)
    textbox(slide, "No narration during generation. None.",
            0.6, 3.4, 8.8, 0.4, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    textbox(slide, "5 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s6_what_happened(prs):
    """Slide 6 — What Just Happened"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    textbox(slide, "Claude didn't fill a template. It made decisions.",
            0.5, 0.15, 9, 0.6, size=22, bold=True, color=TEXT)
    divider(slide, 0.88)

    items = [
        "Generated narrative structure, not just bullet points",
        "Selected what to include and what to cut",
        "Wrote speaker notes so you know what to say",
        "Applied design theme in the same pass",
    ]
    bullets(slide, items, 0.5, 1.0, 9, 2.8, size=17)
    textbox(slide, "6 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s7_architecture(prs):
    """Slide 7 — The Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    textbox(slide, "Built to be trusted.",
            0.5, 0.15, 9, 0.6, size=22, bold=True, color=TEXT)
    divider(slide, 0.88)

    items = [
        "Claude Opus 4.6 with adaptive thinking — the same model used in enterprise deployments",
        "Schema-constrained JSON output — typed, parseable, version-controllable. Same schema every run.",
        "python-pptx rendering — native Office format, no conversion artifacts",
        "Three themes (dark, light, corporate) — designed for real environments, not demos",
    ]
    bullets(slide, items, 0.5, 1.0, 9, 3.5, size=15)
    textbox(slide, "7 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s8_output_options(prs):
    """Slide 8 — Output Options"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    textbox(slide, "Boardroom or browser — same 60 seconds.",
            0.5, 0.15, 9, 0.6, size=22, bold=True, color=TEXT)
    divider(slide, 0.88)

    # Col A
    col_bg_a = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(4.2), Inches(3.6))
    col_bg_a.fill.solid()
    col_bg_a.fill.fore_color.rgb = BG
    col_bg_a.line.color.rgb = DIVIDER

    textbox(slide, "PowerPoint (.pptx)", 0.65, 1.15, 3.9, 0.4,
            size=13, bold=True, color=ACCENT)
    bullets(slide, ["Editable in Office", "Speaker notes intact",
                    "Native formatting", "Enterprise-ready"],
            0.65, 1.6, 3.9, 2.6, size=14)

    # Col B
    col_bg_b = slide.shapes.add_shape(1, Inches(5.2), Inches(1.05), Inches(4.2), Inches(3.6))
    col_bg_b.fill.solid()
    col_bg_b.fill.fore_color.rgb = BG
    col_bg_b.line.color.rgb = DIVIDER

    textbox(slide, "HTML", 5.35, 1.15, 3.9, 0.4,
            size=13, bold=True, color=ACCENT)
    bullets(slide, ["No install required", "Share a link",
                    "Runs in any browser", "Zero dependencies"],
            5.35, 1.6, 3.9, 2.6, size=14)

    textbox(slide, "8 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s9_extension(prs):
    """Slide 9 — The Extension Story"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, SLIDE)
    accent_bar(slide)

    textbox(slide, "The tool you build on, not around.",
            0.5, 0.15, 9, 0.6, size=22, bold=True, color=TEXT)
    divider(slide, 0.88)

    items = [
        "Add a theme in 10 lines of Python — your brand, your colors, available via --theme",
        "Point --slides at any JSON file — your data, your structure, the tool handles layout and notes",
        "Wrap it in a web form: one field, one button, your whole team has access",
        "Schedule it: weekly status deck generated and emailed every Monday, zero human intervention",
    ]
    bullets(slide, items, 0.5, 1.0, 9, 3.5, size=15)
    textbox(slide, "9 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s10_close(prs):
    """Slide 10 — The Only Slide That Matters"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide, ACCENT)

    textbox(slide, "One command. Every presentation.",
            0.6, 0.6, 8.8, 0.7, size=26, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    code_block(slide, 'python generate.py "Your Topic" --theme corporate',
               0.8, 1.7, 8.4, 0.65)
    textbox(slide, "10 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def s11_ask(prs):
    """Slide 11 — The Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    accent_bar(slide, ACCENT)

    textbox(slide, "Run it before you leave this room.",
            0.5, 0.15, 9, 0.6, size=24, bold=True, color=TEXT)
    divider(slide, 0.88)

    commands = [
        "pip install -r requirements.txt",
        'export ANTHROPIC_API_KEY=your-key',
        'python generate.py "Your next meeting"',
    ]
    top = 1.05
    for cmd in commands:
        code_block(slide, cmd, 0.8, top, 8.4, 0.5)
        top += 0.62

    textbox(slide, "github.com/Moodswing9/Headbanger-s-Little-Repository",
            0.5, 3.2, 9, 0.4, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    textbox(slide, "11 / 11", 8.8, 5.2, 1, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    s1_hook(prs)
    s2_real_cost(prs)
    s3_status_quo(prs)
    s4_premise(prs)
    s5_demo(prs)
    s6_what_happened(prs)
    s7_architecture(prs)
    s8_output_options(prs)
    s9_extension(prs)
    s10_close(prs)
    s11_ask(prs)

    out = "AI_Presentation_Generator_Pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}  (11 slides)")


if __name__ == "__main__":
    build()
